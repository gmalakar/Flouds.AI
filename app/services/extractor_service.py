# =============================================================================
# File: extractor_service.py
# Date: 2025-12-23
# Copyright (c) 2024 Goutam Malakar. All rights reserved.
# =============================================================================

import base64
import binascii
import csv
import re
import time
from io import BytesIO, StringIO
from typing import List

import pdfplumber
from bs4 import BeautifulSoup
from docx import Document

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    import openpyxl
except ImportError:
    openpyxl = None

from app.exceptions import InferenceError, ModelNotFoundError
from app.logger import get_logger
from app.models.extracted_file_content import ExtractedFileContent
from app.models.extracted_response import ExtractedResponse
from app.models.file_request import FileRequest
from app.utils.log_sanitizer import sanitize_for_log

logger = get_logger("extractor_service")


class ExtractorService:
    """Service for extraction operations."""

    @staticmethod
    def _normalize_whitespace(text: str) -> str:
        """Normalize tabs/spaces and collapse multiple newlines, similar to the C# helper."""
        if not text:
            return text
        # Replace runs of spaces/tabs with a single space
        normalized = re.sub(r"[ \t]+", " ", text)
        # Replace runs of newlines (any mix of \r/\n) with a single newline
        normalized = re.sub(r"[\r\n]+", "\n", normalized)
        return normalized.strip()

    @staticmethod
    def extract_text(req: FileRequest) -> ExtractedResponse:
        """Main text extraction function."""
        start_time = time.time()
        response = ExtractedResponse(
            success=True,
            message="Text extracted successfully",
            results=[],
            time_taken=0.0,
        )

        try:
            response.results = ExtractorService._extract_local(req)
        except ModelNotFoundError as e:
            logger.warning("Model/library not found: %s", sanitize_for_log(str(e)))
            response.success = False
            response.message = str(e)
        except InferenceError as e:
            logger.error("Extraction error: %s", sanitize_for_log(str(e)))
            response.success = False
            response.message = str(e)
        except Exception as e:
            logger.error(
                "Unexpected extraction error: %s",
                sanitize_for_log(str(e)),
                exc_info=True,
            )
            response.success = False
            response.message = "Error extracting text from file"

        response.time_taken = time.time() - start_time
        return response

    @staticmethod
    def _extract_local(req: FileRequest) -> List[ExtractedFileContent]:
        """Extract text from file based on extension."""
        # Accept either raw bytes or base64-encoded string
        if isinstance(req.file_content, bytes):
            file_bytes = req.file_content
        elif isinstance(req.file_content, str):
            try:
                file_bytes = base64.b64decode(req.file_content)
            except (binascii.Error, ValueError):
                raise InferenceError("Invalid base64 encoded file content")
        else:
            raise InferenceError("file_content must be base64 string or bytes")

        ext = req.extention.lower()

        if ext == "pdf":
            return ExtractorService._extract_pdf(file_bytes)
        elif ext in ["doc", "docx"]:
            return ExtractorService._extract_doc(file_bytes)
        elif ext == "csv":
            return ExtractorService._extract_csv(file_bytes)
        elif ext in ["txt", "md"]:
            return ExtractorService._extract_text(file_bytes)
        elif ext in ["html", "htm"]:
            return ExtractorService._extract_html(file_bytes)
        elif ext in ["ppt", "pptx"]:
            return ExtractorService._extract_ppt(file_bytes)
        elif ext in ["xls", "xlsx"]:
            return ExtractorService._extract_excel(file_bytes)

        raise ModelNotFoundError(f"Unsupported file type: .{ext}")

    @staticmethod
    def _extract_pdf(file_bytes: bytes) -> List[ExtractedFileContent]:
        """Extract text from PDF file."""
        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
            pages = [
                ExtractorService._normalize_whitespace(page.extract_text() or "")
                for page in pdf.pages
            ]
        return [
            ExtractedFileContent(content=page, item_number=i + 1, content_as="pages")
            for i, page in enumerate(pages)
        ]

    @staticmethod
    def _extract_doc(file_bytes: bytes) -> List[ExtractedFileContent]:
        """Extract text from DOC/DOCX file using python-docx."""
        try:
            with BytesIO(file_bytes) as bio:
                doc = Document(bio)
                paragraphs = [
                    ExtractorService._normalize_whitespace(p.text)
                    for p in doc.paragraphs
                    if p.text.strip()
                ]
            return [
                ExtractedFileContent(
                    content=para, item_number=i + 1, content_as="paragraphs"
                )
                for i, para in enumerate(paragraphs)
            ]
        except Exception as e:
            raise InferenceError(
                f"Failed to parse Word document - file may be corrupted or not a valid Word document: {str(e)}"
            )

    # Alias for consistency
    _extract_docx = _extract_doc

    @staticmethod
    def _extract_csv(file_bytes: bytes) -> List[ExtractedFileContent]:
        """Extract text from CSV file."""
        try:
            text = file_bytes.decode("utf-8")
            csv_reader = csv.DictReader(StringIO(text))
            rows = list(csv_reader)
        except UnicodeDecodeError:
            raise InferenceError("CSV file is not valid UTF-8 encoded")
        except Exception as e:
            raise InferenceError(f"Failed to parse CSV file: {str(e)}")

        return [
            ExtractedFileContent(
                content=ExtractorService._normalize_whitespace(str(row)),
                item_number=i + 1,
                content_as="rows",
            )
            for i, row in enumerate(rows)
        ]

    @staticmethod
    def _extract_text(file_bytes: bytes) -> List[ExtractedFileContent]:
        """Extract text from plain text file."""
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            raise InferenceError("File is not valid UTF-8 encoded text")
        normalized = ExtractorService._normalize_whitespace(text)
        return [
            ExtractedFileContent(content=normalized, item_number=1, content_as="text")
        ]

    @staticmethod
    def _extract_html(file_bytes: bytes) -> List[ExtractedFileContent]:
        """Extract text from HTML file."""
        soup = BeautifulSoup(file_bytes, "html.parser")
        return [
            ExtractedFileContent(
                content=ExtractorService._normalize_whitespace(soup.get_text()),
                item_number=1,
                content_as="text",
            )
        ]

    @staticmethod
    def _extract_ppt(file_bytes: bytes) -> List[ExtractedFileContent]:
        """Extract text from PowerPoint file."""
        if not Presentation:
            raise ModelNotFoundError(
                "Install python-pptx to support PowerPoint files: pip install python-pptx"
            )
        try:
            with BytesIO(file_bytes) as bio:
                prs = Presentation(bio)
                slides = []
                for slide in prs.slides:
                    slide_text = [
                        ExtractorService._normalize_whitespace(
                            getattr(shape, "text", "")
                        )
                        for shape in slide.shapes
                    ]
                    slides.append("\n".join(slide_text))
            return [
                ExtractedFileContent(
                    content=slide, item_number=i + 1, content_as="slides"
                )
                for i, slide in enumerate(slides)
            ]
        except Exception as e:
            raise InferenceError(f"Failed to parse PowerPoint file: {str(e)}")

    @staticmethod
    def _extract_excel(file_bytes: bytes) -> List[ExtractedFileContent]:
        """Extract text from Excel file (XLS/XLSX)."""
        if not openpyxl:
            raise ModelNotFoundError(
                "Install openpyxl to support Excel files: pip install openpyxl"
            )

        try:
            with BytesIO(file_bytes) as bio:
                wb = openpyxl.load_workbook(bio, data_only=True)
                sheets = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_data = []
                    for row in ws.iter_rows(values_only=True):
                        if any(cell is not None for cell in row):
                            row_str = ",".join(
                                str(cell) if cell is not None else "" for cell in row
                            )
                            sheet_data.append(
                                ExtractorService._normalize_whitespace(row_str)
                            )
                    sheets.append("\n".join(sheet_data))

            return [
                ExtractedFileContent(
                    content=sheet, item_number=i + 1, content_as="sheets"
                )
                for i, sheet in enumerate(sheets)
            ]
        except Exception as e:
            raise InferenceError(f"Failed to parse Excel file: {str(e)}")
